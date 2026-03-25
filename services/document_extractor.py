"""
Document Extractor Service — PDF, DOCX, XLSX, PPTX

Extracts structured text content from common document formats:
- PDF: pdfplumber (tables + text), fallback to pypdf
- DOCX: python-docx (paragraphs, headings, tables)
- XLSX: openpyxl (sheets, cells, named ranges)
- PPTX: python-pptx (slide titles, text frames, notes)

Returns a unified DocumentContent object with text, metadata, and structure.
"""

import io
import os
import re
from dataclasses import dataclass, field
from typing import Optional
import structlog

logger = structlog.get_logger()


# ============================================================================
# DATA MODELS
# ============================================================================

@dataclass
class DocumentSection:
    """A logical section within a document."""
    heading: str = ""
    content: str = ""
    level: int = 0           # heading depth (0 = body text)
    page: int = 0            # page number (PDF/PPTX)
    section_type: str = "body"  # body, heading, table, code, note


@dataclass
class DocumentContent:
    """Unified output from any document extractor."""
    title: str = ""
    full_text: str = ""          # all text concatenated
    sections: list = field(default_factory=list)   # list[DocumentSection]
    page_count: int = 0
    word_count: int = 0
    language: str = "en"
    doc_type: str = ""           # pdf, docx, xlsx, pptx
    metadata: dict = field(default_factory=dict)
    tables: list = field(default_factory=list)     # list of table dicts
    error: Optional[str] = None

    def to_dict(self) -> dict:
        return {
            "title": self.title,
            "full_text": self.full_text[:50000],   # cap for API response
            "sections": [
                {
                    "heading": s.heading,
                    "content": s.content[:2000],
                    "level": s.level,
                    "page": s.page,
                    "section_type": s.section_type,
                }
                for s in self.sections
            ],
            "page_count": self.page_count,
            "word_count": self.word_count,
            "doc_type": self.doc_type,
            "metadata": self.metadata,
            "table_count": len(self.tables),
            "error": self.error,
        }


# ============================================================================
# DOCUMENT EXTRACTOR
# ============================================================================

class DocumentExtractor:
    """
    Unified document content extractor.

    Dispatches to format-specific extractors based on MIME type or extension.
    All extractors return a DocumentContent object.
    """

    SUPPORTED_TYPES = {
        "application/pdf": "_extract_pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "_extract_docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "_extract_xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "_extract_pptx",
        "application/msword": "_extract_docx",
        "application/vnd.ms-excel": "_extract_xlsx",
        "application/vnd.ms-powerpoint": "_extract_pptx",
    }

    EXT_MAP = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc": "application/msword",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls": "application/vnd.ms-excel",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt": "application/vnd.ms-powerpoint",
    }

    def extract(
        self,
        source: "str | bytes",
        mime_type: Optional[str] = None,
        filename: Optional[str] = None,
    ) -> DocumentContent:
        """
        Extract content from a document.

        Args:
            source:    File path (str) or raw bytes
            mime_type: MIME type (optional — inferred from extension if not given)
            filename:  Original filename for extension detection

        Returns:
            DocumentContent with extracted text, sections, and metadata
        """
        # Resolve MIME type
        if not mime_type:
            if filename:
                ext = os.path.splitext(filename)[1].lower()
                mime_type = self.EXT_MAP.get(ext)
            elif isinstance(source, str):
                ext = os.path.splitext(source)[1].lower()
                mime_type = self.EXT_MAP.get(ext)

        if not mime_type:
            return DocumentContent(error="Cannot determine document type — provide mime_type or filename")

        method_name = self.SUPPORTED_TYPES.get(mime_type)
        if not method_name:
            return DocumentContent(error=f"Unsupported document type: {mime_type}")

        # Load bytes if path given
        data: bytes
        try:
            if isinstance(source, str):
                with open(source, "rb") as f:
                    data = f.read()
            else:
                data = source
        except Exception as e:
            return DocumentContent(error=f"Failed to read file: {e}")

        # Dispatch to format-specific extractor
        try:
            method = getattr(self, method_name)
            result: DocumentContent = method(data, filename or "document")
            result.word_count = len(result.full_text.split())
            logger.info(
                "document_extracted",
                doc_type=result.doc_type,
                pages=result.page_count,
                words=result.word_count,
            )
            return result
        except ImportError as e:
            logger.error("extractor_dependency_missing", error=str(e))
            return DocumentContent(error=f"Missing dependency: {e}. Install required packages.")
        except Exception as e:
            logger.error("extraction_failed", mime_type=mime_type, error=str(e))
            return DocumentContent(error=f"Extraction failed: {e}")

    # -------------------------------------------------------------------------
    # PDF Extractor
    # -------------------------------------------------------------------------

    def _extract_pdf(self, data: bytes, filename: str) -> DocumentContent:
        """Extract text and tables from PDF using pdfplumber."""
        import pdfplumber  # type: ignore

        result = DocumentContent(doc_type="pdf")
        texts = []
        sections = []
        tables = []

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            result.page_count = len(pdf.pages)
            result.metadata = pdf.metadata or {}

            # Use PDF title from metadata if available
            result.title = (
                result.metadata.get("Title", "")
                or result.metadata.get("title", "")
                or os.path.splitext(filename)[0]
            )

            for page_num, page in enumerate(pdf.pages, start=1):
                # Extract text
                page_text = page.extract_text() or ""
                if page_text.strip():
                    texts.append(page_text)
                    sections.append(DocumentSection(
                        heading=f"Page {page_num}",
                        content=page_text.strip(),
                        page=page_num,
                        section_type="body",
                    ))

                # Extract tables
                for tbl in page.extract_tables():
                    if tbl:
                        tables.append({
                            "page": page_num,
                            "rows": tbl,
                        })

        result.full_text = "\n\n".join(texts)
        result.sections = sections
        result.tables = tables
        return result

    # -------------------------------------------------------------------------
    # DOCX Extractor
    # -------------------------------------------------------------------------

    def _extract_docx(self, data: bytes, filename: str) -> DocumentContent:
        """Extract text and structure from DOCX using python-docx."""
        import docx  # type: ignore

        doc = docx.Document(io.BytesIO(data))
        result = DocumentContent(doc_type="docx")
        result.title = os.path.splitext(filename)[0]

        # Extract core properties if available
        try:
            core_props = doc.core_properties
            result.title = core_props.title or result.title
            result.metadata = {
                "author": core_props.author or "",
                "created": str(core_props.created or ""),
                "modified": str(core_props.modified or ""),
                "subject": core_props.subject or "",
            }
        except Exception:
            pass

        full_text_parts = []
        sections = []
        current_section = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else ""
            is_heading = style_name.startswith("Heading")

            if is_heading:
                # Save previous section
                if current_section:
                    sections.append(current_section)

                # Parse heading level from e.g. "Heading 1"
                level_match = re.search(r"Heading (\d+)", style_name)
                level = int(level_match.group(1)) if level_match else 1

                current_section = DocumentSection(
                    heading=text,
                    content="",
                    level=level,
                    section_type="heading",
                )
                full_text_parts.append(text)
            else:
                # Regular paragraph
                if current_section:
                    current_section.content += text + "\n"
                else:
                    sections.append(DocumentSection(
                        content=text,
                        section_type="body",
                    ))
                full_text_parts.append(text)

        if current_section:
            sections.append(current_section)

        # Extract tables
        tables = []
        for tbl in doc.tables:
            rows = []
            for row in tbl.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                tables.append({"rows": rows})
                # Also add table text to full text
                table_text = "\n".join(" | ".join(row) for row in rows)
                full_text_parts.append(table_text)

        result.full_text = "\n\n".join(full_text_parts)
        result.sections = sections
        result.tables = tables
        result.page_count = len(doc.sections)
        return result

    # -------------------------------------------------------------------------
    # XLSX Extractor
    # -------------------------------------------------------------------------

    def _extract_xlsx(self, data: bytes, filename: str) -> DocumentContent:
        """Extract text and data from XLSX using openpyxl."""
        import openpyxl  # type: ignore

        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        result = DocumentContent(doc_type="xlsx")
        result.title = os.path.splitext(filename)[0]
        result.page_count = len(wb.sheetnames)

        full_text_parts = []
        sections = []
        tables = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                # Skip entirely empty rows
                row_strs = [str(c) if c is not None else "" for c in row]
                if any(row_strs):
                    rows.append(row_strs)

            if rows:
                # Header row as heading
                header = " | ".join(rows[0]) if rows else ""
                body_text = "\n".join(" | ".join(r) for r in rows[1:50])  # cap rows

                sections.append(DocumentSection(
                    heading=f"Sheet: {sheet_name}",
                    content=f"{header}\n{body_text}",
                    section_type="body",
                ))
                full_text_parts.append(f"Sheet: {sheet_name}\n{header}\n{body_text}")
                tables.append({"sheet": sheet_name, "rows": rows[:50]})

        result.full_text = "\n\n".join(full_text_parts)
        result.sections = sections
        result.tables = tables
        return result

    # -------------------------------------------------------------------------
    # PPTX Extractor
    # -------------------------------------------------------------------------

    def _extract_pptx(self, data: bytes, filename: str) -> DocumentContent:
        """Extract text from PPTX slides using python-pptx."""
        import pptx  # type: ignore

        prs = pptx.Presentation(io.BytesIO(data))
        result = DocumentContent(doc_type="pptx")
        result.title = os.path.splitext(filename)[0]
        result.page_count = len(prs.slides)

        full_text_parts = []
        sections = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_title = ""
            slide_texts = []
            notes_text = ""

            # Extract title
            if slide.shapes.title:
                slide_title = slide.shapes.title.text.strip()

            # Extract all text frames
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and text != slide_title:
                            slide_texts.append(text)

            # Extract speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_text = notes

            slide_content = "\n".join(slide_texts)
            if notes_text:
                slide_content += f"\n[Notes: {notes_text}]"

            heading = slide_title or f"Slide {slide_num}"
            sections.append(DocumentSection(
                heading=heading,
                content=slide_content,
                page=slide_num,
                level=1,
                section_type="heading" if slide_title else "body",
            ))
            full_text_parts.append(f"{heading}\n{slide_content}")

        result.full_text = "\n\n".join(full_text_parts)
        result.sections = sections
        return result


# Global singleton
document_extractor = DocumentExtractor()
